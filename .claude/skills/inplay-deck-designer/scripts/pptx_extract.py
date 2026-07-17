#!/usr/bin/env python3
"""
pptx_extract.py, InPlay deck designer ingestion (Phase 2, Step 2).

Extracts a source PowerPoint deck into everything the transformation needs:
  - per-slide text as Markdown (title, body in shape order, tables)
  - speaker notes per slide
  - every picture, saved per slide
  - manifest.json summarising what was found

Usage:
    python3 pptx_extract.py <deck.pptx> <output_dir>

Output layout:
    <output_dir>/content.md              slide-by-slide text and tables
    <output_dir>/images/slide-NN-*.png   pictures per slide
    <output_dir>/manifest.json           slides, shape/word/image counts, flags

Notes:
  - Grouped shapes are walked recursively so nested text is not lost.
  - Slides that are mostly pictures (little text) are flagged
    "needs_visual_read"; inspect the extracted images visually.
  - House style: long-dash characters (U+2014, U+2013) in extracted text are
    normalised to commas. Nothing else about the wording is changed.
  - Keynote files are not supported; ask Brett to export to PPTX or PDF first.
"""
import json
import os
import re
import sys

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    sys.exit("python-pptx is required. Install with: pip3 install python-pptx")

# U+2014 (em dash) and U+2013 (en dash) are banned by house style
DASH_RE = re.compile("\\s*[\\u2014\\u2013]\\s*")


def clean(text: str) -> str:
    text = DASH_RE.sub(", ", text)
    return re.sub(r"[ \t]+", " ", text).strip()


def walk_shapes(shapes):
    """Yield leaf shapes, recursing into groups."""
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from walk_shapes(shape.shapes)
        else:
            yield shape


def shape_markdown(shape) -> str:
    if shape.has_table:
        rows = []
        for r, row in enumerate(shape.table.rows):
            cells = [clean(c.text) for c in row.cells]
            rows.append("| " + " | ".join(cells) + " |")
            if r == 0:
                rows.append("|" + "---|" * len(cells))
        return "\n".join(rows)
    if shape.has_text_frame:
        lines = []
        for para in shape.text_frame.paragraphs:
            text = clean("".join(run.text for run in para.runs))
            if not text:
                continue
            indent = "  " * (para.level or 0)
            lines.append(f"{indent}- {text}")
        return "\n".join(lines)
    return ""


def main() -> None:
    if len(sys.argv) != 3:
        sys.exit(__doc__)
    pptx_path, out_dir = sys.argv[1], sys.argv[2]
    images_dir = os.path.join(out_dir, "images")
    os.makedirs(images_dir, exist_ok=True)

    prs = Presentation(pptx_path)
    md_parts = [f"# Source deck: {os.path.basename(pptx_path)}\n"]
    manifest = {"source": os.path.abspath(pptx_path), "slides": []}

    for i, slide in enumerate(prs.slides, start=1):
        title = ""
        if slide.shapes.title is not None and slide.shapes.title.text.strip():
            title = clean(slide.shapes.title.text)

        body_parts = []
        img_count = 0
        for shape in walk_shapes(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                    ext = image.ext or "png"
                    path = os.path.join(images_dir, f"slide-{i:02d}-{img_count + 1}.{ext}")
                    with open(path, "wb") as f:
                        f.write(image.blob)
                    img_count += 1
                except Exception:
                    continue
                continue
            if shape == slide.shapes.title:
                continue
            md = shape_markdown(shape)
            if md:
                body_parts.append(md)

        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
            notes = clean(slide.notes_slide.notes_text_frame.text)

        body = "\n\n".join(body_parts)
        section = [f"\n---\n\n# Slide {i}"]
        if title:
            section.append(f"## {title}")
        section.append(body or "(no body text)")
        if img_count:
            section.append(f"*{img_count} picture(s) extracted to images/slide-{i:02d}-\\**")
        if notes:
            section.append(f"> Speaker notes: {notes}")
        md_parts.append("\n\n".join(section))

        words = len((title + " " + body).split())
        manifest["slides"].append({
            "slide": i,
            "title": title,
            "words": words,
            "pictures": img_count,
            "has_notes": bool(notes),
            "needs_visual_read": img_count > 0 and words < 15,
        })

    with open(os.path.join(out_dir, "content.md"), "w") as f:
        f.write("\n".join(md_parts) + "\n")
    with open(os.path.join(out_dir, "manifest.json"), "w") as f:
        json.dump(manifest, f, indent=2)

    flagged = sum(1 for s in manifest["slides"] if s["needs_visual_read"])
    print(f"Extracted {len(manifest['slides'])} slides to {out_dir}"
          f" ({flagged} slide(s) flagged for visual read)")


if __name__ == "__main__":
    main()
